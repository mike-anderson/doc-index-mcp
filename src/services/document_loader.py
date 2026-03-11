"""
Document loading service for extracting text from various file formats.

Supports:
- Plain text (.txt)
- Markdown (.md)
- PDF (.pdf) - using pdfplumber
- Word documents (.docx) - using python-docx
- PowerPoint presentations (.pptx) - using python-pptx
- Excel spreadsheets (.xlsx, .xls) - using openpyxl
"""

import os
from dataclasses import dataclass, field
from typing import Optional

from .chunker import Boundary, BoundaryType


SUPPORTED_EXTENSIONS = {
    '.txt': 'text',
    '.md': 'markdown',
    '.markdown': 'markdown',
    '.pdf': 'pdf',
    '.docx': 'docx',
    '.pptx': 'pptx',
    '.xlsx': 'xlsx',
    '.xls': 'xlsx',
}


@dataclass
class LoadedDocument:
    """Result of loading a document."""
    content: str
    file_type: str
    metadata: dict
    structure: Optional[dict] = field(default=None)
    boundaries: list[Boundary] = field(default_factory=list)  # Pre-detected by native loaders


async def load_document(file_path: str) -> LoadedDocument:
    """
    Load and extract text from a document.

    Args:
        file_path: Path to the document file

    Returns:
        LoadedDocument with extracted content and metadata

    Raises:
        FileNotFoundError: If file doesn't exist
        ValueError: If file type is not supported
    """
    if not os.path.exists(file_path):
        raise FileNotFoundError(f"File not found: {file_path}")

    ext = os.path.splitext(file_path)[1].lower()

    if ext not in SUPPORTED_EXTENSIONS:
        raise ValueError(f"Unsupported file type: {ext}. Supported: {', '.join(SUPPORTED_EXTENSIONS.keys())}")

    if ext in ('.txt', '.md', '.markdown'):
        return await _load_text_file(file_path, ext)
    elif ext == '.pdf':
        return await _load_pdf_file(file_path)
    elif ext == '.docx':
        return await _load_docx_file(file_path)
    elif ext == '.pptx':
        return await _load_pptx_file(file_path)
    elif ext in ('.xlsx', '.xls'):
        return await _load_xlsx_file(file_path)
    else:
        raise ValueError(f"Unsupported file type: {ext}")


async def _load_text_file(file_path: str, ext: str) -> LoadedDocument:
    """Load a plain text or markdown file."""
    with open(file_path, 'r', encoding='utf-8') as f:
        content = f.read()

    file_type = 'markdown' if ext in ('.md', '.markdown') else 'text'

    return LoadedDocument(
        content=content,
        file_type=file_type,
        metadata={
            "file_path": file_path,
            "file_size": os.path.getsize(file_path),
        }
    )


async def _load_pdf_file(file_path: str) -> LoadedDocument:
    """Load a PDF file using pdfplumber with native page boundary detection."""
    import pdfplumber

    pages_content = []
    boundaries = []
    total_pages = 0
    offset = 0

    with pdfplumber.open(file_path) as pdf:
        total_pages = len(pdf.pages)

        for i, page in enumerate(pdf.pages, start=1):
            text = page.extract_text() or ""
            if text.strip():
                page_text = f"[Page {i}]\n{text}"

                # Create boundary from native page knowledge
                boundaries.append(Boundary(
                    type=BoundaryType.PAGE,
                    level=4,
                    id=f"page:{i}",
                    title=str(i),
                    start_offset=offset,
                ))

                pages_content.append(page_text)
                offset += len(page_text) + 2  # +2 for \n\n separator

    content = "\n\n".join(pages_content)

    return LoadedDocument(
        content=content,
        file_type='pdf',
        metadata={
            "file_path": file_path,
            "file_size": os.path.getsize(file_path),
            "total_pages": total_pages,
        },
        structure={
            "pages": total_pages,
        },
        boundaries=boundaries,
    )


async def _load_docx_file(file_path: str) -> LoadedDocument:
    """
    Load a Word document using python-docx.

    Extracts:
    - Paragraphs as text blocks
    - Tables converted to markdown tables
    - Headings marked with # notation for boundary detection
    - Page breaks detected from native <w:br w:type="page"/> elements
    """
    from docx import Document
    from lxml import etree

    doc = Document(file_path)
    content_parts = []
    boundaries = []
    heading_count = 0
    table_count = 0
    page_number = 1
    offset = 0

    # XML namespace for detecting page breaks
    w_ns = '{http://schemas.openxmlformats.org/wordprocessingml/2006/main}'

    def _has_page_break(element) -> bool:
        """Check if a paragraph element contains a page break."""
        # Check for <w:br w:type="page"/> in runs
        for br in element.iter(f'{w_ns}br'):
            if br.get(f'{w_ns}type') == 'page':
                return True
        # Check for page break before in paragraph properties
        pPr = element.find(f'{w_ns}pPr')
        if pPr is not None:
            pageBreakBefore = pPr.find(f'{w_ns}pageBreakBefore')
            if pageBreakBefore is not None:
                val = pageBreakBefore.get(f'{w_ns}val')
                if val is None or val.lower() not in ('false', '0', 'off'):
                    return True
        # Check for section breaks (which also start new pages)
        if pPr is not None:
            sectPr = pPr.find(f'{w_ns}sectPr')
            if sectPr is not None:
                return True
        return False

    for element in doc.element.body:
        if element.tag.endswith('p'):
            # Check for page break before processing content
            if _has_page_break(element):
                page_number += 1
                page_marker = f"[Page {page_number}]"
                boundaries.append(Boundary(
                    type=BoundaryType.PAGE,
                    level=4,
                    id=f"page:{page_number}",
                    title=str(page_number),
                    start_offset=offset,
                ))
                content_parts.append(page_marker)
                offset += len(page_marker) + 2

            # Find matching paragraph
            for para in doc.paragraphs:
                if para._element is element:
                    text = para.text.strip()
                    if not text:
                        continue

                    # Detect heading styles
                    if para.style and para.style.name and para.style.name.startswith('Heading'):
                        heading_count += 1
                        level_str = para.style.name.replace('Heading', '').strip()
                        try:
                            level = int(level_str) if level_str.isdigit() else 1
                        except ValueError:
                            level = 1
                        level = min(level, 6)
                        part = f"{'#' * level} {text}"
                    else:
                        part = text

                    content_parts.append(part)
                    offset += len(part) + 2
                    break

        elif element.tag.endswith('tbl'):
            for table in doc.tables:
                if table._tbl is element:
                    table_count += 1
                    md_table = _docx_table_to_markdown(table)
                    if md_table:
                        content_parts.append(md_table)
                        offset += len(md_table) + 2
                    break

    content = "\n\n".join(content_parts)

    return LoadedDocument(
        content=content,
        file_type='docx',
        metadata={
            "file_path": file_path,
            "file_size": os.path.getsize(file_path),
            "total_pages": page_number,
        },
        structure={
            "headings": heading_count,
            "tables": table_count,
            "pages": page_number,
        },
        boundaries=boundaries,
    )


def _docx_table_to_markdown(table) -> str:
    """Convert a python-docx table to markdown format."""
    rows = []
    for row in table.rows:
        cells = [cell.text.strip().replace('\n', ' ') for cell in row.cells]
        rows.append(cells)

    if not rows:
        return ""

    return _rows_to_markdown_table(rows)


async def _load_pptx_file(file_path: str) -> LoadedDocument:
    """
    Load a PowerPoint presentation using python-pptx.

    Extracts:
    - Each slide as [Slide N: Title] boundary
    - Slide text content
    - Speaker notes marked with [Notes:] section
    - Tables converted to markdown
    """
    from pptx import Presentation

    prs = Presentation(file_path)
    content_parts = []
    boundaries = []
    slide_count = len(prs.slides)
    has_notes = False
    offset = 0

    for i, slide in enumerate(prs.slides, start=1):
        slide_parts = []

        # Get slide title
        title = _get_slide_title(slide)
        slide_parts.append(f"[Slide {i}: {title}]")

        # Create boundary from native slide knowledge
        boundaries.append(Boundary(
            type=BoundaryType.SLIDE,
            level=1,
            id=f"slide:{i}",
            title=title,
            start_offset=offset,
        ))

        # Extract text from shapes
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = _extract_text_from_text_frame(shape.text_frame)
                if text and text != title:  # Avoid duplicating title
                    slide_parts.append(text)

            if shape.has_table:
                md_table = _pptx_table_to_markdown(shape.table)
                if md_table:
                    slide_parts.append(md_table)

        # Speaker notes
        if slide.has_notes_slide:
            notes_frame = slide.notes_slide.notes_text_frame
            if notes_frame and notes_frame.text.strip():
                has_notes = True
                slide_parts.append(f"\n[Notes:]\n{notes_frame.text.strip()}")

        slide_text = "\n".join(slide_parts)
        content_parts.append(slide_text)
        offset += len(slide_text) + 2  # +2 for \n\n separator

    content = "\n\n".join(content_parts)

    return LoadedDocument(
        content=content,
        file_type='pptx',
        metadata={
            "file_path": file_path,
            "file_size": os.path.getsize(file_path),
            "total_slides": slide_count,
        },
        structure={
            "slides": slide_count,
            "has_notes": has_notes,
        },
        boundaries=boundaries,
    )


def _get_slide_title(slide) -> str:
    """Extract the title from a PowerPoint slide."""
    if slide.shapes.title:
        return slide.shapes.title.text.strip() or "Untitled"

    # Fallback: look for title placeholder
    for shape in slide.shapes:
        if shape.is_placeholder:
            if hasattr(shape, 'placeholder_format'):
                # Check for title placeholder types
                if shape.placeholder_format.type in (1, 3):  # TITLE or CENTER_TITLE
                    if shape.has_text_frame:
                        return shape.text_frame.text.strip() or "Untitled"

    return "Untitled"


def _extract_text_from_text_frame(text_frame) -> str:
    """Extract text from a text frame, preserving paragraph structure."""
    paragraphs = []
    for para in text_frame.paragraphs:
        text = para.text.strip()
        if text:
            paragraphs.append(text)
    return "\n".join(paragraphs)


def _pptx_table_to_markdown(table) -> str:
    """Convert a python-pptx table to markdown format."""
    rows = []
    for row in table.rows:
        cells = []
        for cell in row.cells:
            # Extract text from cell
            text = cell.text.strip().replace('\n', ' ')
            cells.append(text)
        rows.append(cells)

    if not rows:
        return ""

    return _rows_to_markdown_table(rows)


async def _load_xlsx_file(
    file_path: str,
    rows_per_group: int = 50,
) -> LoadedDocument:
    """
    Load an Excel spreadsheet using openpyxl with native sheet and row boundaries.

    Extracts:
    - Each sheet as a SHEET boundary
    - Row groups as ROW_GROUP boundaries within each sheet
    - Data converted to markdown tables per row group
    """
    from openpyxl import load_workbook

    # data_only=True gets calculated values instead of formulas
    wb = load_workbook(file_path, data_only=True)
    content_parts = []
    boundaries = []
    total_rows = 0
    offset = 0
    sheet_names = list(wb.sheetnames)

    for sheet_idx, sheet_name in enumerate(sheet_names, start=1):
        sheet = wb[sheet_name]

        # Sheet boundary
        sheet_marker = f"[Sheet: {sheet_name}]"
        boundaries.append(Boundary(
            type=BoundaryType.SHEET,
            level=1,
            id=f"sheet:{sheet_idx}",
            title=sheet_name,
            start_offset=offset,
        ))
        content_parts.append(sheet_marker)
        offset += len(sheet_marker) + 2

        # Collect all non-empty rows
        all_rows = []
        for row in sheet.iter_rows(values_only=True):
            if any(cell is not None for cell in row):
                all_rows.append(row)
                total_rows += 1

        if not all_rows:
            empty_text = "*Empty sheet*"
            content_parts.append(empty_text)
            offset += len(empty_text) + 2
            continue

        # Extract header row
        header_row = all_rows[0]
        data_rows = all_rows[1:]

        if not data_rows:
            # Only a header row
            md_table = _rows_to_markdown_table(all_rows, max_rows=500)
            content_parts.append(md_table)
            offset += len(md_table) + 2
            continue

        # Split data rows into groups for row-level boundaries
        for group_start in range(0, len(data_rows), rows_per_group):
            group_end = min(group_start + rows_per_group, len(data_rows))
            group_rows = data_rows[group_start:group_end]

            # Row numbers are 1-indexed (excluding header)
            row_start_num = group_start + 2  # +2: 1-indexed, skip header
            row_end_num = group_end + 1

            row_marker = f"[Rows {row_start_num}-{row_end_num}]"
            boundaries.append(Boundary(
                type=BoundaryType.ROW_GROUP,
                level=2,
                id=f"sheet:{sheet_idx}:rows:{row_start_num}-{row_end_num}",
                title=f"{sheet_name} rows {row_start_num}-{row_end_num}",
                start_offset=offset,
                parent_id=f"sheet:{sheet_idx}",
            ))

            # Build markdown table with header repeated for each group
            table_rows = [header_row] + list(group_rows)
            md_table = _rows_to_markdown_table(table_rows, max_rows=rows_per_group + 1)

            group_text = f"{row_marker}\n{md_table}"
            content_parts.append(group_text)
            offset += len(group_text) + 2

    wb.close()
    content = "\n\n".join(content_parts)

    return LoadedDocument(
        content=content,
        file_type='xlsx',
        metadata={
            "file_path": file_path,
            "file_size": os.path.getsize(file_path),
        },
        structure={
            "sheets": sheet_names,
            "total_rows": total_rows,
        },
        boundaries=boundaries,
    )


def _rows_to_markdown_table(rows: list, max_rows: int = 500) -> str:
    """
    Convert rows to markdown table format.

    Args:
        rows: List of tuples/lists representing table rows
        max_rows: Maximum number of data rows to include (default 500)

    Returns:
        Markdown formatted table string
    """
    if not rows:
        return ""

    # Determine max columns across all rows
    max_cols = max(len(row) for row in rows)

    # Normalize all rows to have the same number of columns
    normalized_rows = []
    for row in rows:
        normalized = list(row) + [None] * (max_cols - len(row))
        normalized_rows.append(normalized)

    # Use first row as headers
    headers = []
    for i, h in enumerate(normalized_rows[0]):
        if h is not None:
            headers.append(str(h).replace('|', '\\|').replace('\n', ' '))
        else:
            headers.append(f"Col{i + 1}")

    lines = [
        "| " + " | ".join(headers) + " |",
        "| " + " | ".join(["---"] * len(headers)) + " |",
    ]

    # Add data rows (skip header row)
    data_rows = normalized_rows[1:max_rows + 1]
    for row in data_rows:
        cells = []
        for c in row:
            if c is not None:
                cell_str = str(c).replace('|', '\\|').replace('\n', ' ')
                cells.append(cell_str)
            else:
                cells.append("")
        lines.append("| " + " | ".join(cells) + " |")

    if len(normalized_rows) > max_rows + 1:
        remaining = len(normalized_rows) - max_rows - 1
        lines.append(f"\n*... {remaining} more rows truncated*")

    return "\n".join(lines)


def get_source_name(file_path: str, custom_name: Optional[str] = None) -> str:
    """
    Generate a source name from a file path.

    Args:
        file_path: Path to the file
        custom_name: Optional custom name to use instead

    Returns:
        Source name string
    """
    if custom_name:
        # Sanitize custom name
        return custom_name.replace(' ', '_').replace('/', '_')

    # Use filename without extension
    basename = os.path.basename(file_path)
    name, _ = os.path.splitext(basename)
    return name.replace(' ', '_')


def get_supported_extensions() -> list[str]:
    """Return list of supported file extensions."""
    return list(SUPPORTED_EXTENSIONS.keys())
