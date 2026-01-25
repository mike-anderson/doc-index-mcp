"""
Table extraction service for extracting structured table data from documents.

Supports:
- Excel spreadsheets (.xlsx, .xls) - each sheet is a table
- Word documents (.docx) - embedded tables
- PowerPoint presentations (.pptx) - tables in slides
- PDF files (.pdf) - detected tables via pdfplumber
"""

import csv
import io
import os
from dataclasses import dataclass
from typing import Optional


@dataclass
class TableInfo:
    """Information about a table in a document."""
    index: int
    location: str  # e.g., "Sheet1", "Page 3", "Slide 2"
    headers: list[str]
    row_count: int
    col_count: int


@dataclass
class ExtractedTable:
    """Extracted table data."""
    index: int
    location: str
    headers: list[str]
    rows: list[list]  # List of rows, each row is a list of cell values
    row_count: int
    col_count: int


async def list_tables(file_path: str) -> list[TableInfo]:
    """
    List all tables in a document.

    Args:
        file_path: Path to the document file

    Returns:
        List of TableInfo objects describing each table

    Raises:
        FileNotFoundError: If file doesn't exist
        ValueError: If file type is not supported
    """
    if not os.path.exists(file_path):
        raise FileNotFoundError(f"File not found: {file_path}")

    ext = os.path.splitext(file_path)[1].lower()

    if ext in ('.xlsx', '.xls'):
        return await _list_xlsx_tables(file_path)
    elif ext == '.docx':
        return await _list_docx_tables(file_path)
    elif ext == '.pptx':
        return await _list_pptx_tables(file_path)
    elif ext == '.pdf':
        return await _list_pdf_tables(file_path)
    else:
        raise ValueError(f"Unsupported file type for table extraction: {ext}")


async def extract_table(
    file_path: str,
    table_index: int,
    max_rows: Optional[int] = None,
) -> ExtractedTable:
    """
    Extract a specific table from a document.

    Args:
        file_path: Path to the document file
        table_index: Index of the table to extract (0-based)
        max_rows: Maximum number of data rows to extract (None = all)

    Returns:
        ExtractedTable with the table data

    Raises:
        FileNotFoundError: If file doesn't exist
        ValueError: If file type not supported or table_index out of range
    """
    if not os.path.exists(file_path):
        raise FileNotFoundError(f"File not found: {file_path}")

    ext = os.path.splitext(file_path)[1].lower()

    if ext in ('.xlsx', '.xls'):
        return await _extract_xlsx_table(file_path, table_index, max_rows)
    elif ext == '.docx':
        return await _extract_docx_table(file_path, table_index, max_rows)
    elif ext == '.pptx':
        return await _extract_pptx_table(file_path, table_index, max_rows)
    elif ext == '.pdf':
        return await _extract_pdf_table(file_path, table_index, max_rows)
    else:
        raise ValueError(f"Unsupported file type for table extraction: {ext}")


def table_to_csv(table: ExtractedTable, include_headers: bool = True) -> str:
    """
    Convert an extracted table to CSV format.

    Args:
        table: ExtractedTable to convert
        include_headers: Whether to include header row

    Returns:
        CSV formatted string
    """
    output = io.StringIO()
    writer = csv.writer(output)

    if include_headers and table.headers:
        writer.writerow(table.headers)

    for row in table.rows:
        # Convert all values to strings, handling None
        str_row = [str(cell) if cell is not None else "" for cell in row]
        writer.writerow(str_row)

    return output.getvalue()


# Excel handlers

async def _list_xlsx_tables(file_path: str) -> list[TableInfo]:
    """List tables in an Excel file (each sheet = one table)."""
    from openpyxl import load_workbook

    wb = load_workbook(file_path, data_only=True, read_only=True)
    tables = []

    for idx, sheet_name in enumerate(wb.sheetnames):
        sheet = wb[sheet_name]

        # Get dimensions
        rows = list(sheet.iter_rows(values_only=True))
        # Filter empty rows
        rows = [r for r in rows if any(cell is not None for cell in r)]

        if not rows:
            # Empty sheet, still list it but with 0 rows
            tables.append(TableInfo(
                index=idx,
                location=sheet_name,
                headers=[],
                row_count=0,
                col_count=0,
            ))
            continue

        # First row as headers
        headers = [str(h) if h is not None else f"Col{i+1}" for i, h in enumerate(rows[0])]
        col_count = len(headers)
        row_count = len(rows) - 1  # Exclude header row

        tables.append(TableInfo(
            index=idx,
            location=sheet_name,
            headers=headers,
            row_count=row_count,
            col_count=col_count,
        ))

    wb.close()
    return tables


async def _extract_xlsx_table(
    file_path: str,
    table_index: int,
    max_rows: Optional[int],
) -> ExtractedTable:
    """Extract a specific sheet from an Excel file."""
    from openpyxl import load_workbook

    wb = load_workbook(file_path, data_only=True, read_only=True)

    if table_index < 0 or table_index >= len(wb.sheetnames):
        wb.close()
        raise ValueError(f"Table index {table_index} out of range. File has {len(wb.sheetnames)} sheets.")

    sheet_name = wb.sheetnames[table_index]
    sheet = wb[sheet_name]

    # Get all rows
    all_rows = list(sheet.iter_rows(values_only=True))
    # Filter empty rows
    all_rows = [r for r in all_rows if any(cell is not None for cell in r)]

    wb.close()

    if not all_rows:
        return ExtractedTable(
            index=table_index,
            location=sheet_name,
            headers=[],
            rows=[],
            row_count=0,
            col_count=0,
        )

    # Normalize column count
    max_cols = max(len(row) for row in all_rows)
    normalized_rows = []
    for row in all_rows:
        normalized = list(row) + [None] * (max_cols - len(row))
        normalized_rows.append(normalized)

    # First row as headers
    headers = [str(h) if h is not None else f"Col{i+1}" for i, h in enumerate(normalized_rows[0])]

    # Data rows (skip header)
    data_rows = normalized_rows[1:]
    if max_rows is not None:
        data_rows = data_rows[:max_rows]

    return ExtractedTable(
        index=table_index,
        location=sheet_name,
        headers=headers,
        rows=data_rows,
        row_count=len(data_rows),
        col_count=len(headers),
    )


# Word handlers

async def _list_docx_tables(file_path: str) -> list[TableInfo]:
    """List tables in a Word document."""
    from docx import Document

    doc = Document(file_path)
    tables = []

    for idx, table in enumerate(doc.tables):
        rows = []
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            rows.append(cells)

        if not rows:
            tables.append(TableInfo(
                index=idx,
                location=f"Table {idx + 1}",
                headers=[],
                row_count=0,
                col_count=0,
            ))
            continue

        headers = rows[0] if rows else []
        row_count = len(rows) - 1 if rows else 0

        tables.append(TableInfo(
            index=idx,
            location=f"Table {idx + 1}",
            headers=headers,
            row_count=row_count,
            col_count=len(headers),
        ))

    return tables


async def _extract_docx_table(
    file_path: str,
    table_index: int,
    max_rows: Optional[int],
) -> ExtractedTable:
    """Extract a specific table from a Word document."""
    from docx import Document

    doc = Document(file_path)

    if table_index < 0 or table_index >= len(doc.tables):
        raise ValueError(f"Table index {table_index} out of range. Document has {len(doc.tables)} tables.")

    table = doc.tables[table_index]

    rows = []
    for row in table.rows:
        cells = [cell.text.strip() for cell in row.cells]
        rows.append(cells)

    if not rows:
        return ExtractedTable(
            index=table_index,
            location=f"Table {table_index + 1}",
            headers=[],
            rows=[],
            row_count=0,
            col_count=0,
        )

    headers = rows[0]
    data_rows = rows[1:]
    if max_rows is not None:
        data_rows = data_rows[:max_rows]

    return ExtractedTable(
        index=table_index,
        location=f"Table {table_index + 1}",
        headers=headers,
        rows=data_rows,
        row_count=len(data_rows),
        col_count=len(headers),
    )


# PowerPoint handlers

async def _list_pptx_tables(file_path: str) -> list[TableInfo]:
    """List tables in a PowerPoint presentation."""
    from pptx import Presentation

    prs = Presentation(file_path)
    tables = []
    table_idx = 0

    for slide_num, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            if shape.has_table:
                table = shape.table
                rows = []
                for row in table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    rows.append(cells)

                headers = rows[0] if rows else []
                row_count = len(rows) - 1 if rows else 0

                tables.append(TableInfo(
                    index=table_idx,
                    location=f"Slide {slide_num}",
                    headers=headers,
                    row_count=row_count,
                    col_count=len(headers),
                ))
                table_idx += 1

    return tables


async def _extract_pptx_table(
    file_path: str,
    table_index: int,
    max_rows: Optional[int],
) -> ExtractedTable:
    """Extract a specific table from a PowerPoint presentation."""
    from pptx import Presentation

    prs = Presentation(file_path)
    table_idx = 0

    for slide_num, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            if shape.has_table:
                if table_idx == table_index:
                    table = shape.table
                    rows = []
                    for row in table.rows:
                        cells = [cell.text.strip() for cell in row.cells]
                        rows.append(cells)

                    if not rows:
                        return ExtractedTable(
                            index=table_index,
                            location=f"Slide {slide_num}",
                            headers=[],
                            rows=[],
                            row_count=0,
                            col_count=0,
                        )

                    headers = rows[0]
                    data_rows = rows[1:]
                    if max_rows is not None:
                        data_rows = data_rows[:max_rows]

                    return ExtractedTable(
                        index=table_index,
                        location=f"Slide {slide_num}",
                        headers=headers,
                        rows=data_rows,
                        row_count=len(data_rows),
                        col_count=len(headers),
                    )
                table_idx += 1

    raise ValueError(f"Table index {table_index} out of range. Presentation has {table_idx} tables.")


# PDF handlers

async def _list_pdf_tables(file_path: str) -> list[TableInfo]:
    """List tables in a PDF file using pdfplumber."""
    import pdfplumber

    tables = []
    table_idx = 0

    with pdfplumber.open(file_path) as pdf:
        for page_num, page in enumerate(pdf.pages, start=1):
            page_tables = page.extract_tables()

            for page_table in page_tables:
                if not page_table or not page_table[0]:
                    continue

                # First row as headers
                headers = [str(h) if h else f"Col{i+1}" for i, h in enumerate(page_table[0])]
                row_count = len(page_table) - 1

                tables.append(TableInfo(
                    index=table_idx,
                    location=f"Page {page_num}",
                    headers=headers,
                    row_count=row_count,
                    col_count=len(headers),
                ))
                table_idx += 1

    return tables


async def _extract_pdf_table(
    file_path: str,
    table_index: int,
    max_rows: Optional[int],
) -> ExtractedTable:
    """Extract a specific table from a PDF file."""
    import pdfplumber

    table_idx = 0

    with pdfplumber.open(file_path) as pdf:
        for page_num, page in enumerate(pdf.pages, start=1):
            page_tables = page.extract_tables()

            for page_table in page_tables:
                if not page_table or not page_table[0]:
                    continue

                if table_idx == table_index:
                    # First row as headers
                    headers = [str(h) if h else f"Col{i+1}" for i, h in enumerate(page_table[0])]

                    # Data rows
                    data_rows = page_table[1:]
                    if max_rows is not None:
                        data_rows = data_rows[:max_rows]

                    # Normalize None values
                    normalized_rows = []
                    for row in data_rows:
                        normalized = [cell if cell is not None else "" for cell in row]
                        normalized_rows.append(normalized)

                    return ExtractedTable(
                        index=table_index,
                        location=f"Page {page_num}",
                        headers=headers,
                        rows=normalized_rows,
                        row_count=len(normalized_rows),
                        col_count=len(headers),
                    )
                table_idx += 1

    raise ValueError(f"Table index {table_index} out of range. PDF has {table_idx} tables.")


def get_supported_table_formats() -> list[str]:
    """Return list of file extensions that support table extraction."""
    return ['.xlsx', '.xls', '.docx', '.pptx', '.pdf']
