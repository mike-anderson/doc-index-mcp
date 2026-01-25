"""
Tests for document loading service.

Tests format-specific loaders for DOCX, PPTX, XLSX and the unified load_document function.
"""

import pytest
import os
import tempfile
from io import BytesIO

from ..services.document_loader import (
    LoadedDocument,
    load_document,
    get_source_name,
    get_supported_extensions,
    _rows_to_markdown_table,
    SUPPORTED_EXTENSIONS,
)


class TestSupportedExtensions:
    """Test extension mapping and helper functions."""

    def test_supported_extensions_dict(self):
        """Should have all expected formats mapped."""
        assert '.txt' in SUPPORTED_EXTENSIONS
        assert '.md' in SUPPORTED_EXTENSIONS
        assert '.pdf' in SUPPORTED_EXTENSIONS
        assert '.docx' in SUPPORTED_EXTENSIONS
        assert '.pptx' in SUPPORTED_EXTENSIONS
        assert '.xlsx' in SUPPORTED_EXTENSIONS
        assert '.xls' in SUPPORTED_EXTENSIONS

    def test_get_supported_extensions(self):
        """Should return list of supported extensions."""
        exts = get_supported_extensions()
        assert isinstance(exts, list)
        assert '.docx' in exts
        assert '.pptx' in exts
        assert '.xlsx' in exts


class TestMarkdownTableHelper:
    """Test the markdown table conversion helper."""

    def test_simple_table(self):
        """Should convert simple rows to markdown table."""
        rows = [
            ("Name", "Age", "City"),
            ("Alice", 30, "NYC"),
            ("Bob", 25, "LA"),
        ]
        result = _rows_to_markdown_table(rows)

        assert "| Name | Age | City |" in result
        assert "| --- | --- | --- |" in result
        assert "| Alice | 30 | NYC |" in result
        assert "| Bob | 25 | LA |" in result

    def test_empty_rows(self):
        """Should return empty string for empty rows."""
        result = _rows_to_markdown_table([])
        assert result == ""

    def test_none_values(self):
        """Should handle None values in cells."""
        rows = [
            ("Header1", "Header2"),
            ("Value", None),
        ]
        result = _rows_to_markdown_table(rows)

        assert "| Header1 | Header2 |" in result
        assert "| Value |  |" in result

    def test_pipe_escape(self):
        """Should escape pipe characters in cell content."""
        rows = [
            ("Data",),
            ("a|b",),
        ]
        result = _rows_to_markdown_table(rows)

        # Pipe should be escaped
        assert "a\\|b" in result

    def test_newline_replacement(self):
        """Should replace newlines with spaces."""
        rows = [
            ("Text",),
            ("line1\nline2",),
        ]
        result = _rows_to_markdown_table(rows)

        # Newlines should become spaces
        assert "line1 line2" in result
        assert "\n" not in result.split("\n")[2]  # In data row

    def test_truncation(self):
        """Should truncate large tables with message."""
        rows = [("Header",)] + [(f"Row{i}",) for i in range(600)]
        result = _rows_to_markdown_table(rows, max_rows=500)

        # Should have truncation message
        assert "more rows truncated" in result
        # Should not have all rows
        assert "Row550" not in result

    def test_auto_headers_for_empty(self):
        """Should generate Col1, Col2 headers for None values."""
        rows = [
            (None, "Header2", None),
            ("a", "b", "c"),
        ]
        result = _rows_to_markdown_table(rows)

        assert "Col1" in result
        assert "Header2" in result
        assert "Col3" in result

    def test_uneven_rows(self):
        """Should handle rows with different column counts."""
        rows = [
            ("A", "B", "C"),
            ("1", "2"),  # Missing column
            ("X", "Y", "Z", "Extra"),  # Extra column ignored by normalization
        ]
        result = _rows_to_markdown_table(rows)

        # Should create valid table
        assert "| A | B | C |" in result


class TestLoadDocument:
    """Test the unified load_document function."""

    def test_file_not_found(self):
        """Should raise FileNotFoundError for missing files."""
        with pytest.raises(FileNotFoundError):
            import asyncio
            asyncio.run(load_document("/nonexistent/file.txt"))

    def test_unsupported_format(self, tmp_path):
        """Should raise ValueError for unsupported formats."""
        # Create a file with unsupported extension
        test_file = tmp_path / "test.xyz"
        test_file.write_text("content")

        with pytest.raises(ValueError, match="Unsupported file type"):
            import asyncio
            asyncio.run(load_document(str(test_file)))

    def test_load_text_file(self, tmp_path):
        """Should load plain text files."""
        test_file = tmp_path / "test.txt"
        test_file.write_text("Hello, World!")

        import asyncio
        doc = asyncio.run(load_document(str(test_file)))

        assert doc.content == "Hello, World!"
        assert doc.file_type == "text"
        assert doc.metadata["file_path"] == str(test_file)

    def test_load_markdown_file(self, tmp_path):
        """Should load markdown files."""
        test_file = tmp_path / "test.md"
        test_file.write_text("# Heading\n\nParagraph text.")

        import asyncio
        doc = asyncio.run(load_document(str(test_file)))

        assert "# Heading" in doc.content
        assert doc.file_type == "markdown"


class TestGetSourceName:
    """Test source name generation."""

    def test_from_path(self):
        """Should extract filename without extension."""
        name = get_source_name("/path/to/document.pdf")
        assert name == "document"

    def test_custom_name(self):
        """Should use custom name when provided."""
        name = get_source_name("/path/to/file.txt", custom_name="my_source")
        assert name == "my_source"

    def test_sanitize_spaces(self):
        """Should replace spaces with underscores."""
        name = get_source_name("/path/to/my document.pdf")
        assert name == "my_document"

    def test_sanitize_custom_name(self):
        """Should sanitize custom names."""
        name = get_source_name("/path/to/file.txt", custom_name="my source/name")
        assert name == "my_source_name"


class TestDocxLoader:
    """Test Word document loading."""

    @pytest.fixture
    def simple_docx(self, tmp_path):
        """Create a simple DOCX file for testing."""
        from docx import Document

        doc = Document()
        doc.add_heading("Test Document", level=1)
        doc.add_paragraph("This is a test paragraph.")
        doc.add_heading("Section 1", level=2)
        doc.add_paragraph("Content in section 1.")

        path = tmp_path / "test.docx"
        doc.save(str(path))
        return str(path)

    @pytest.fixture
    def docx_with_table(self, tmp_path):
        """Create a DOCX with a table."""
        from docx import Document

        doc = Document()
        doc.add_paragraph("Document with table:")

        table = doc.add_table(rows=3, cols=2)
        table.cell(0, 0).text = "Name"
        table.cell(0, 1).text = "Value"
        table.cell(1, 0).text = "Item A"
        table.cell(1, 1).text = "100"
        table.cell(2, 0).text = "Item B"
        table.cell(2, 1).text = "200"

        path = tmp_path / "table.docx"
        doc.save(str(path))
        return str(path)

    def test_load_simple_docx(self, simple_docx):
        """Should extract text from simple DOCX."""
        import asyncio
        doc = asyncio.run(load_document(simple_docx))

        assert doc.file_type == "docx"
        assert "Test Document" in doc.content
        assert "test paragraph" in doc.content

    def test_docx_headings_to_markdown(self, simple_docx):
        """Should convert headings to markdown format."""
        import asyncio
        doc = asyncio.run(load_document(simple_docx))

        # Heading 1 should become #
        assert "# Test Document" in doc.content
        # Heading 2 should become ##
        assert "## Section 1" in doc.content

    def test_docx_with_tables(self, docx_with_table):
        """Should convert tables to markdown."""
        import asyncio
        doc = asyncio.run(load_document(docx_with_table))

        # Should have markdown table
        assert "| Name | Value |" in doc.content
        assert "| Item A | 100 |" in doc.content

    def test_docx_structure_metadata(self, simple_docx):
        """Should include structure info in metadata."""
        import asyncio
        doc = asyncio.run(load_document(simple_docx))

        assert doc.structure is not None
        assert "headings" in doc.structure
        assert doc.structure["headings"] >= 2  # At least 2 headings


class TestPptxLoader:
    """Test PowerPoint loading."""

    @pytest.fixture
    def simple_pptx(self, tmp_path):
        """Create a simple PPTX for testing."""
        from pptx import Presentation
        from pptx.util import Inches

        prs = Presentation()

        # Add title slide
        slide_layout = prs.slide_layouts[0]  # Title slide layout
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = "Presentation Title"
        subtitle = slide.placeholders[1]
        subtitle.text = "Subtitle text"

        # Add content slide
        slide_layout = prs.slide_layouts[1]  # Title and content
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = "Slide 2 Title"
        body = slide.placeholders[1]
        tf = body.text_frame
        tf.text = "First bullet point"
        p = tf.add_paragraph()
        p.text = "Second bullet point"

        path = tmp_path / "test.pptx"
        prs.save(str(path))
        return str(path)

    @pytest.fixture
    def pptx_with_notes(self, tmp_path):
        """Create a PPTX with speaker notes."""
        from pptx import Presentation

        prs = Presentation()
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = "Slide with Notes"

        # Add notes
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = "These are the speaker notes."

        path = tmp_path / "notes.pptx"
        prs.save(str(path))
        return str(path)

    def test_load_simple_pptx(self, simple_pptx):
        """Should extract text from PPTX."""
        import asyncio
        doc = asyncio.run(load_document(simple_pptx))

        assert doc.file_type == "pptx"
        assert "Presentation Title" in doc.content
        assert "Slide 2 Title" in doc.content

    def test_pptx_slide_markers(self, simple_pptx):
        """Should include slide boundary markers."""
        import asyncio
        doc = asyncio.run(load_document(simple_pptx))

        assert "[Slide 1:" in doc.content
        assert "[Slide 2:" in doc.content

    def test_pptx_with_notes(self, pptx_with_notes):
        """Should extract speaker notes."""
        import asyncio
        doc = asyncio.run(load_document(pptx_with_notes))

        assert "[Notes:]" in doc.content
        assert "speaker notes" in doc.content
        assert doc.structure["has_notes"] is True

    def test_pptx_structure_metadata(self, simple_pptx):
        """Should include slide count in structure."""
        import asyncio
        doc = asyncio.run(load_document(simple_pptx))

        assert doc.structure is not None
        assert doc.structure["slides"] == 2


class TestXlsxLoader:
    """Test Excel spreadsheet loading."""

    @pytest.fixture
    def simple_xlsx(self, tmp_path):
        """Create a simple XLSX for testing."""
        from openpyxl import Workbook

        wb = Workbook()
        ws = wb.active
        ws.title = "Data"

        # Add header row
        ws["A1"] = "Name"
        ws["B1"] = "Age"
        ws["C1"] = "City"

        # Add data rows
        ws["A2"] = "Alice"
        ws["B2"] = 30
        ws["C2"] = "NYC"

        ws["A3"] = "Bob"
        ws["B3"] = 25
        ws["C3"] = "LA"

        path = tmp_path / "test.xlsx"
        wb.save(str(path))
        return str(path)

    @pytest.fixture
    def multisheet_xlsx(self, tmp_path):
        """Create an XLSX with multiple sheets."""
        from openpyxl import Workbook

        wb = Workbook()

        # First sheet
        ws1 = wb.active
        ws1.title = "Sheet1"
        ws1["A1"] = "Header1"
        ws1["A2"] = "Value1"

        # Second sheet
        ws2 = wb.create_sheet("Sheet2")
        ws2["A1"] = "Header2"
        ws2["A2"] = "Value2"

        path = tmp_path / "multi.xlsx"
        wb.save(str(path))
        return str(path)

    def test_load_simple_xlsx(self, simple_xlsx):
        """Should extract data from XLSX as markdown table."""
        import asyncio
        doc = asyncio.run(load_document(simple_xlsx))

        assert doc.file_type == "xlsx"
        assert "| Name | Age | City |" in doc.content
        assert "| Alice | 30 | NYC |" in doc.content

    def test_xlsx_sheet_markers(self, simple_xlsx):
        """Should include sheet boundary markers."""
        import asyncio
        doc = asyncio.run(load_document(simple_xlsx))

        assert "[Sheet: Data]" in doc.content

    def test_xlsx_multisheet(self, multisheet_xlsx):
        """Should load all sheets."""
        import asyncio
        doc = asyncio.run(load_document(multisheet_xlsx))

        assert "[Sheet: Sheet1]" in doc.content
        assert "[Sheet: Sheet2]" in doc.content
        assert "Header1" in doc.content
        assert "Header2" in doc.content

    def test_xlsx_structure_metadata(self, multisheet_xlsx):
        """Should include sheet names in structure."""
        import asyncio
        doc = asyncio.run(load_document(multisheet_xlsx))

        assert doc.structure is not None
        assert "sheets" in doc.structure
        assert "Sheet1" in doc.structure["sheets"]
        assert "Sheet2" in doc.structure["sheets"]

    def test_xlsx_empty_cells(self, tmp_path):
        """Should handle empty cells gracefully."""
        from openpyxl import Workbook

        wb = Workbook()
        ws = wb.active
        ws["A1"] = "A"
        ws["C1"] = "C"  # B1 is empty
        ws["A2"] = "1"
        ws["B2"] = "2"
        ws["C2"] = "3"

        path = tmp_path / "sparse.xlsx"
        wb.save(str(path))

        import asyncio
        doc = asyncio.run(load_document(str(path)))

        # Should handle the missing B1 cell
        assert "| A |" in doc.content


class TestLoadedDocument:
    """Test LoadedDocument dataclass."""

    def test_create_with_defaults(self):
        """Should create with optional structure as None."""
        doc = LoadedDocument(
            content="test",
            file_type="text",
            metadata={"key": "value"},
        )
        assert doc.structure is None

    def test_create_with_structure(self):
        """Should accept structure dict."""
        doc = LoadedDocument(
            content="test",
            file_type="docx",
            metadata={},
            structure={"headings": 5, "tables": 2},
        )
        assert doc.structure["headings"] == 5
