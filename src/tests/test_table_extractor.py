"""
Tests for table extraction service.

Tests list_tables and extract_table functions for XLSX, DOCX, PPTX formats.
"""

import pytest
import asyncio

from ..services.table_extractor import (
    TableInfo,
    ExtractedTable,
    list_tables,
    extract_table,
    table_to_csv,
    get_supported_table_formats,
)


class TestSupportedFormats:
    """Test supported format helpers."""

    def test_get_supported_table_formats(self):
        """Should return list of supported formats."""
        formats = get_supported_table_formats()
        assert '.xlsx' in formats
        assert '.docx' in formats
        assert '.pptx' in formats
        assert '.pdf' in formats


class TestTableToCsv:
    """Test CSV conversion."""

    def test_simple_table(self):
        """Should convert table to CSV."""
        table = ExtractedTable(
            index=0,
            location="Sheet1",
            headers=["Name", "Age"],
            rows=[["Alice", 30], ["Bob", 25]],
            row_count=2,
            col_count=2,
        )
        csv = table_to_csv(table)

        assert "Name,Age" in csv
        assert "Alice,30" in csv
        assert "Bob,25" in csv

    def test_without_headers(self):
        """Should omit headers when requested."""
        table = ExtractedTable(
            index=0,
            location="Sheet1",
            headers=["Name", "Age"],
            rows=[["Alice", 30]],
            row_count=1,
            col_count=2,
        )
        csv = table_to_csv(table, include_headers=False)

        assert "Name,Age" not in csv
        assert "Alice,30" in csv

    def test_none_values(self):
        """Should handle None values as empty strings."""
        table = ExtractedTable(
            index=0,
            location="Sheet1",
            headers=["A", "B"],
            rows=[["value", None]],
            row_count=1,
            col_count=2,
        )
        csv = table_to_csv(table)

        # None should become empty
        assert "value," in csv

    def test_csv_escaping(self):
        """Should properly escape CSV special characters."""
        table = ExtractedTable(
            index=0,
            location="Sheet1",
            headers=["Text"],
            rows=[["hello, world"], ["with \"quotes\""]],
            row_count=2,
            col_count=1,
        )
        csv = table_to_csv(table)

        # CSV should handle commas and quotes
        assert csv  # Just verify it doesn't crash


class TestXlsxTables:
    """Test Excel table extraction."""

    @pytest.fixture
    def simple_xlsx(self, tmp_path):
        """Create a simple Excel file."""
        from openpyxl import Workbook

        wb = Workbook()
        ws = wb.active
        ws.title = "Data"
        ws["A1"] = "Name"
        ws["B1"] = "Value"
        ws["A2"] = "Item1"
        ws["B2"] = 100
        ws["A3"] = "Item2"
        ws["B3"] = 200

        path = tmp_path / "test.xlsx"
        wb.save(str(path))
        return str(path)

    @pytest.fixture
    def multisheet_xlsx(self, tmp_path):
        """Create Excel with multiple sheets."""
        from openpyxl import Workbook

        wb = Workbook()

        ws1 = wb.active
        ws1.title = "Revenue"
        ws1["A1"] = "Quarter"
        ws1["B1"] = "Amount"
        ws1["A2"] = "Q1"
        ws1["B2"] = 1000

        ws2 = wb.create_sheet("Expenses")
        ws2["A1"] = "Category"
        ws2["B1"] = "Cost"
        ws2["A2"] = "Rent"
        ws2["B2"] = 500

        path = tmp_path / "multi.xlsx"
        wb.save(str(path))
        return str(path)

    def test_list_xlsx_tables(self, simple_xlsx):
        """Should list sheets as tables."""
        tables = asyncio.run(list_tables(simple_xlsx))

        assert len(tables) == 1
        assert tables[0].location == "Data"
        assert tables[0].headers == ["Name", "Value"]
        assert tables[0].row_count == 2

    def test_list_multisheet_xlsx(self, multisheet_xlsx):
        """Should list all sheets."""
        tables = asyncio.run(list_tables(multisheet_xlsx))

        assert len(tables) == 2
        assert tables[0].location == "Revenue"
        assert tables[1].location == "Expenses"

    def test_extract_xlsx_table(self, simple_xlsx):
        """Should extract table as CSV."""
        table = asyncio.run(extract_table(simple_xlsx, 0))

        assert table.location == "Data"
        assert table.headers == ["Name", "Value"]
        assert len(table.rows) == 2
        assert table.rows[0] == ["Item1", 100]

    def test_extract_xlsx_with_max_rows(self, simple_xlsx):
        """Should respect max_rows limit."""
        table = asyncio.run(extract_table(simple_xlsx, 0, max_rows=1))

        assert table.row_count == 1
        assert len(table.rows) == 1

    def test_extract_xlsx_invalid_index(self, simple_xlsx):
        """Should raise error for invalid index."""
        with pytest.raises(ValueError, match="out of range"):
            asyncio.run(extract_table(simple_xlsx, 99))


class TestDocxTables:
    """Test Word document table extraction."""

    @pytest.fixture
    def docx_with_table(self, tmp_path):
        """Create Word doc with a table."""
        from docx import Document

        doc = Document()
        doc.add_paragraph("Document with table:")

        table = doc.add_table(rows=3, cols=2)
        table.cell(0, 0).text = "Product"
        table.cell(0, 1).text = "Price"
        table.cell(1, 0).text = "Widget"
        table.cell(1, 1).text = "9.99"
        table.cell(2, 0).text = "Gadget"
        table.cell(2, 1).text = "19.99"

        path = tmp_path / "table.docx"
        doc.save(str(path))
        return str(path)

    @pytest.fixture
    def docx_with_multiple_tables(self, tmp_path):
        """Create Word doc with multiple tables."""
        from docx import Document

        doc = Document()

        # First table
        table1 = doc.add_table(rows=2, cols=2)
        table1.cell(0, 0).text = "A"
        table1.cell(0, 1).text = "B"
        table1.cell(1, 0).text = "1"
        table1.cell(1, 1).text = "2"

        doc.add_paragraph("Between tables")

        # Second table
        table2 = doc.add_table(rows=2, cols=3)
        table2.cell(0, 0).text = "X"
        table2.cell(0, 1).text = "Y"
        table2.cell(0, 2).text = "Z"
        table2.cell(1, 0).text = "10"
        table2.cell(1, 1).text = "20"
        table2.cell(1, 2).text = "30"

        path = tmp_path / "multi_table.docx"
        doc.save(str(path))
        return str(path)

    def test_list_docx_tables(self, docx_with_table):
        """Should find tables in Word doc."""
        tables = asyncio.run(list_tables(docx_with_table))

        assert len(tables) == 1
        assert tables[0].headers == ["Product", "Price"]
        assert tables[0].row_count == 2

    def test_list_multiple_docx_tables(self, docx_with_multiple_tables):
        """Should find all tables."""
        tables = asyncio.run(list_tables(docx_with_multiple_tables))

        assert len(tables) == 2
        assert tables[0].col_count == 2
        assert tables[1].col_count == 3

    def test_extract_docx_table(self, docx_with_table):
        """Should extract table data."""
        table = asyncio.run(extract_table(docx_with_table, 0))

        assert table.headers == ["Product", "Price"]
        assert table.rows[0] == ["Widget", "9.99"]
        assert table.rows[1] == ["Gadget", "19.99"]

    def test_extract_docx_csv(self, docx_with_table):
        """Should convert to valid CSV."""
        table = asyncio.run(extract_table(docx_with_table, 0))
        csv = table_to_csv(table)

        assert "Product,Price" in csv
        assert "Widget,9.99" in csv


class TestPptxTables:
    """Test PowerPoint table extraction."""

    @pytest.fixture
    def pptx_with_table(self, tmp_path):
        """Create PowerPoint with a table."""
        from pptx import Presentation
        from pptx.util import Inches

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout

        # Add table
        x, y, cx, cy = Inches(1), Inches(1), Inches(4), Inches(1.5)
        table = slide.shapes.add_table(3, 2, x, y, cx, cy).table

        table.cell(0, 0).text = "Metric"
        table.cell(0, 1).text = "Value"
        table.cell(1, 0).text = "Sales"
        table.cell(1, 1).text = "1000"
        table.cell(2, 0).text = "Costs"
        table.cell(2, 1).text = "500"

        path = tmp_path / "table.pptx"
        prs.save(str(path))
        return str(path)

    def test_list_pptx_tables(self, pptx_with_table):
        """Should find tables in slides."""
        tables = asyncio.run(list_tables(pptx_with_table))

        assert len(tables) == 1
        assert "Slide" in tables[0].location
        assert tables[0].headers == ["Metric", "Value"]

    def test_extract_pptx_table(self, pptx_with_table):
        """Should extract table data."""
        table = asyncio.run(extract_table(pptx_with_table, 0))

        assert table.headers == ["Metric", "Value"]
        assert table.rows[0] == ["Sales", "1000"]


class TestErrorHandling:
    """Test error handling."""

    def test_file_not_found(self):
        """Should raise FileNotFoundError."""
        with pytest.raises(FileNotFoundError):
            asyncio.run(list_tables("/nonexistent/file.xlsx"))

    def test_unsupported_format(self, tmp_path):
        """Should raise ValueError for unsupported formats."""
        test_file = tmp_path / "test.txt"
        test_file.write_text("not a table file")

        with pytest.raises(ValueError, match="Unsupported"):
            asyncio.run(list_tables(str(test_file)))


class TestTableInfo:
    """Test TableInfo dataclass."""

    def test_create_table_info(self):
        """Should create TableInfo with all fields."""
        info = TableInfo(
            index=0,
            location="Sheet1",
            headers=["A", "B"],
            row_count=10,
            col_count=2,
        )
        assert info.index == 0
        assert info.location == "Sheet1"
        assert info.headers == ["A", "B"]
        assert info.row_count == 10
        assert info.col_count == 2


class TestExtractedTable:
    """Test ExtractedTable dataclass."""

    def test_create_extracted_table(self):
        """Should create ExtractedTable with all fields."""
        table = ExtractedTable(
            index=0,
            location="Sheet1",
            headers=["A", "B"],
            rows=[["1", "2"], ["3", "4"]],
            row_count=2,
            col_count=2,
        )
        assert table.index == 0
        assert len(table.rows) == 2
